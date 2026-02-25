#!/usr/bin/env python3
"""Generate a formal 11-slide PayTrace presentation for Techno-verse Hackathon."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design system ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Deep Navy
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)       # Card Navy
ACCENT_PRI = RGBColor(0x10, 0xB9, 0x81)    # Professional Green
ACCENT_SEC = RGBColor(0x3B, 0x82, 0xF6)    # Professional Blue
TEXT_MAIN = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_SEC = RGBColor(0x94, 0xA3, 0xB8)      # Slate Gray

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_title(slide, text, subtext=None):
    # Top bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.4), Inches(0.15), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_PRI
    shape.line.fill.background()
    
    # Title
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(10), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.color.rgb = TEXT_MAIN
    p.font.bold = True
    p.font.name = 'Arial'

    if subtext:
        tx2 = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(10), Inches(0.5))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = subtext
        p2.font.size = Pt(18)
        p2.font.color.rgb = ACCENT_SEC
        p2.font.name = 'Arial'

def add_card(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    return shape

def add_bullet(slide, x, y, w, h, title, points, color=TEXT_MAIN):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = ACCENT_PRI
        p.font.bold = True
        p.space_after = Pt(10)
        
    for pt in points:
        p = tf.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(16)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0

# ════════════════════════════════════════════
# SLIDE 1: TITLE & TEAM
# ════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)

# Hackathon Badge
badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(1.5), Inches(4), Inches(0.6))
badge.fill.solid()
badge.fill.fore_color.rgb = ACCENT_SEC
badge.line.fill.background()
tf = badge.text_frame
tf.paragraphs[0].text = "Techno-verse Hackathon"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.bold = True

add_bullet(s1, Inches(0), Inches(2.5), W, Inches(1.5), "", ["PayTrace"], color=TEXT_MAIN)
s1.shapes[-1].text_frame.paragraphs[1].font.size = Pt(64)
s1.shapes[-1].text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
s1.shapes[-1].text_frame.paragraphs[1].font.bold = True

add_bullet(s1, Inches(0), Inches(3.8), W, Inches(1), "", ["Smart UPI Expense Tracking with AI-Powered Intelligence"], color=TEXT_SEC)
s1.shapes[-1].text_frame.paragraphs[1].font.size = Pt(24)
s1.shapes[-1].text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

# Team
tbox = add_card(s1, Inches(3.6), Inches(5.0), Inches(6), Inches(2.0))
add_bullet(s1, Inches(3.8), Inches(5.1), Inches(5.6), Inches(1.8), "Team Members", 
           ["Danush K", "Jeeva G", "Kavin Soorya S", "Arivuchezhian", "Barath J"])
s1.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ════════════════════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
add_title(s2, "Problem Statement", "The gap in financial tracking")

add_card(s2, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.5))

add_card(s2, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.5))
add_bullet(s2, Inches(0.8), Inches(2.3), Inches(11.5), Inches(3.0), "", [
    "In an ecosystem processing over 12 billion UPI transactions monthly, users suffer from 'Financial Blindness'. Payments are fragmented across multiple apps (GPay, PhonePe, Paytm) and manual tracking is tedious and error-prone.",
    "",
    "Existing solutions are either unintelligent—offering only raw transaction lists—or privacy-invasive, requiring cloud uploads. There is a critical lack of a unified, on-device system that can autonomously track spending and provide intelligent behavioral insights without compromising user privacy."
])

# ════════════════════════════════════════════
# SLIDE 3: PROPOSED SOLUTION
# ════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
add_title(s3, "Proposed Solution", "PayTrace: Privacy-First On-Device Intelligence")

# 3 Pillars
cols = [
    ("Track", "Automatically detects ALL UPI transactions (Sent + Received) by scanning bank SMS. Works for every UPI app."),
    ("Analyze", "On-device ML engine autonomously discovers spending patterns (Time clusters, Recurring payments)."),
    ("Act", "Converts raw data into personalized, actionable insights (Savings goals, Budget alerts).")
]

for i, (title, text) in enumerate(cols):
    x = Inches(0.5 + i * 4.1)
    add_card(s3, x, Inches(2.0), Inches(3.8), Inches(2.5))
    add_bullet(s3, x + Inches(0.2), Inches(2.1), Inches(3.4), Inches(2.3), title, [text])
    s3.shapes[-1].text_frame.paragraphs[1].font.size = Pt(14)

add_card(s3, Inches(0.5), Inches(5.0), Inches(12.0), Inches(1.8))
add_bullet(s3, Inches(0.7), Inches(5.1), Inches(11.5), Inches(1.5), "Key Differentiators", [
    "Universal Coverage: Tracks all apps via SMS.",
    "100% On-Device: No cloud servers. Zero data leak.",
    "Autonomous ML: Learns patterns without manual rules."
])

# ════════════════════════════════════════════
# SLIDE 4: SYSTEM ARCHITECTURE
# ════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4)
add_title(s4, "System Architecture", "6-Layer Processing Pipeline")

layers = [
    "LAYER 1: Data Ingestion (SMS, QR, Contacts)",
    "LAYER 2: Processing Pipeline (Parse, Dedup, Categorize)",
    "LAYER 3: Local Database (SQLite + Drift)",
    "LAYER 4: ML Pattern Engine (Clustering, Anomaly)",
    "LAYER 5: Insight Generation (NLG, Advisories)",
    "LAYER 6: User Interface (Dashboard, Alerts)"
]

y_start = 1.8
for i, layer in enumerate(layers):
    add_card(s4, Inches(3.5), Inches(y_start + i * 0.8), Inches(6.3), Inches(0.6))
    tx = s4.shapes.add_textbox(Inches(3.5), Inches(y_start + i * 0.8 + 0.1), Inches(6.3), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = layer
    p.font.color.rgb = TEXT_MAIN
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True

    if i < 5:
        ar = s4.shapes.add_textbox(Inches(6.5), Inches(y_start + i * 0.8 + 0.6), Inches(0.5), Inches(0.2))
        ar.text_frame.paragraphs[0].text = "⬇"
        ar.text_frame.paragraphs[0].font.color.rgb = TEXT_SEC
        ar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ════════════════════════════════════════════
# SLIDE 5: DATA INGESTION & PROCESSING
# ════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5)
add_title(s5, "Data Ingestion & Processing", "Ensuring Data Integrity")

add_card(s5, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5))
add_bullet(s5, Inches(0.7), Inches(2.0), Inches(5.4), Inches(4), "SMS Sync Engine", [
    "Trigger: Runs on app startup.",
    "Pipeline: Read SMS → Filter Junk → Parse → Dedup → Insert.",
    "Junk Filter: Removes OTPs, Promos, Balance alerts."
])

add_card(s5, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
add_bullet(s5, Inches(7.0), Inches(2.0), Inches(5.4), Inches(4), "Smart Deduplication", [
    "Layer 1: Exact UPI Reference Logic.",
    "Layer 2: Amount + Timestamp fuzzy match (±5 min).",
    "Layer 3: Re-import guard to prevent duplicates."
])

# ════════════════════════════════════════════
# SLIDE 6: CORE MODULES
# ════════════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6)
add_title(s6, "Core Implementation Modules", "Storage & Classification")

add_card(s6, Inches(0.5), Inches(1.8), Inches(12.3), Inches(2.0))
add_bullet(s6, Inches(0.7), Inches(2.0), Inches(11.8), Inches(1.6), "Name Resolution Module", [
    "Extracts phone number from UPI ID (e.g. 98765...@ybl).",
    "Matches against Device Contacts to resolve names.",
    "Fallback Hierarchy: Contacts → User Edits → SMS Body → Bank Code."
])

add_card(s6, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.0))
add_bullet(s6, Inches(0.7), Inches(4.4), Inches(11.8), Inches(1.6), "Auto-Categorization Module", [
    "Rule-based engine with 50+ merchant patterns.",
    "Classifies 'Swiggy' as Food, 'Uber' as Transport.",
    "Automatically tags Credits as 'Income'."
])

# ════════════════════════════════════════════
# SLIDE 7: ML PATTERN ANALYSIS (TEMPORAL)
# ════════════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7)
add_title(s7, "ML Analysis: Temporal & Financial", "Simple Explanation")

# Pattern 1
add_card(s7, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5))
add_bullet(s7, Inches(0.7), Inches(2.0), Inches(5.4), Inches(4.0), "Time-based Clustering", [
    "Algorithm: K-Means Clustering.",
    "How it works: Groups spending by hour of day.",
    "Insight: Detects 'Late Night' clusters usually associated with food delivery or impulse buying."
])

# Pattern 2
add_card(s7, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
add_bullet(s7, Inches(7.0), Inches(2.0), Inches(5.4), Inches(4.0), "Anomaly Detection", [
    "Algorithm: Isolation Forest / Z-Score.",
    "How it works: Flags transactions deviating >2 standard deviations from mean.",
    "Insight: Alerts usage on unusually large purchases."
])

# ════════════════════════════════════════════
# SLIDE 8: ML PATTERN ANALYSIS (BEHAVIORAL)
# ════════════════════════════════════════════
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8)
add_title(s8, "ML Analysis: Behavioral", "Simple Explanation")

# Pattern 3
add_card(s8, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5))
add_bullet(s8, Inches(0.7), Inches(2.0), Inches(5.4), Inches(4.0), "Recurring Payment Detection", [
    "Algorithm: Autocorrelation / FFT.",
    "How it works: Analyzes time intervals between identical payees.",
    "Insight: Identifies subscriptions (Netflix, Rent, SIPs) automatically."
])

# Pattern 4
add_card(s8, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
add_bullet(s8, Inches(7.0), Inches(2.0), Inches(5.4), Inches(4.0), "Merchant Clustering", [
    "Algorithm: DBSCAN.",
    "How it works: Groups payees based on frequency and amount.",
    "Insight: Segments merchants into 'Daily Essentials' vs 'Luxury/Splurge'."
])

# ════════════════════════════════════════════
# SLIDE 9: ACTIONABLE INSIGHTS
# ════════════════════════════════════════════
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9)
add_title(s9, "Actionable Delivery", "From Data to Decision")

add_card(s9, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.5))
add_bullet(s9, Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.0), "User-Facing Features", [
    "Smart Dashboard: Visualizes 'Sent vs Received' in real-time.",
    "Budget Alerts: Warnings at 75% and 90% of set limits.",
    "Natural Language Insights: 'You spend ₹2,000 on weekends'.",
    "Financial Health Score: A 0-100 score based on savings discipline.",
    "Spending Trends: Week-over-week comparison charts."
])

# ════════════════════════════════════════════
# SLIDE 10: IMPACT & FEASIBILITY
# ════════════════════════════════════════════
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10)
add_title(s10, "Impact & Feasibility", "Why this works")

add_card(s10, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.5))
add_bullet(s10, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.0), "Impact", [
    "Financial Discipline: Helps users save 10-15% by identifying leaks.",
    "Privacy Assurance: Zero data upload builds user trust.",
    "Time Saving: Eliminates manual entry."
])

add_card(s10, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.5))
add_bullet(s10, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.0), "Scalability & Feasibility", [
    "Zero Server Cost: Client-side architecture.",
    "Universal: Works with any bank SMS format.",
    "Tech Stack: Flutter + SQLite + TensorFlow Lite is production-ready."
])

# ════════════════════════════════════════════
# SLIDE 11: FUTURE SCOPE
# ════════════════════════════════════════════
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s11)
add_title(s11, "Future Scope", "Roadmap")

add_card(s11, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.5))
add_bullet(s11, Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.0), "", [
    "Multi-Language Support: Parsing SMS in Hindi, Tamil, Telugu.",
    "Generative AI Chatbot: 'How much did I spend on coffee?'.",
    "Family Tracking: Shared expense management.",
    "Encrypted Backup: Optional cloud sync with AES-256.",
    "iOS Support: Using specialized SMS filtering APIs."
])

# Save
output_path = os.path.expanduser("~/Desktop/PayTrace_Hackathon_Formal.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
